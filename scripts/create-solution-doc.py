"""
AssetGuardian AI — Solution Document PowerPoint Generator
Run: pip install python-pptx && python create-solution-doc.py
Output: AssetGuardian_Solution_Document.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Colors
DARK_BG = RGBColor(0x0B, 0x12, 0x20)
ACCENT_BLUE = RGBColor(0x4F, 0x8C, 0xFF)
ACCENT_GREEN = RGBColor(0x6E, 0xE7, 0xB7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x93, 0xA1, 0xBD)
PANEL = RGBColor(0x12, 0x1A, 0x2B)


def add_slide(title, subtitle=None, layout_idx=5):
    """Add a blank slide with dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    return slide


def add_title_box(slide, text, top=0.4, left=0.5, width=12, size=32, color=WHITE, bold=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    return txBox


def add_text_box(slide, text, top, left=0.5, width=12, size=16, color=MUTED, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        p.space_after = Pt(6)
    return txBox


def add_bullet_box(slide, items, top, left=0.5, width=11, size=14, color=MUTED):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  •  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return txBox


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("Title")
add_title_box(slide, "AssetGuardian AI", top=2.5, left=1, size=44, color=ACCENT_GREEN)
add_text_box(slide, "Enterprise Asset Inspection Platform\nAI-Powered Identity Verification, Damage Detection, Fraud Prevention & Lifecycle Management", top=3.5, left=1, size=18, color=MUTED)
add_text_box(slide, "Solution Architecture Document", top=5.0, left=1, size=14, color=ACCENT_BLUE)
add_text_box(slide, "NCS Technology Services  |  Confidential", top=6.5, left=1, size=11, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: Executive Summary
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("Executive Summary")
add_title_box(slide, "Executive Summary", top=0.3, size=28)
add_text_box(slide, "AssetGuardian AI automates enterprise IT asset inspections using multi-agent AI orchestration.\nIt replaces manual, error-prone processes with an intelligent pipeline that verifies identity,\ndetects physical damage, prevents fraud, and makes lifecycle decisions — all in a single pass.", top=1.2, size=15, color=WHITE)
add_bullet_box(slide, [
    "Reduces inspection time from 45 minutes to under 60 seconds",
    "Eliminates human bias and inconsistency in damage assessment",
    "Detects fraud (stock photos, reused images, manipulated submissions) automatically",
    "Auto-identifies device type and brand from photos — no manual selection needed",
    "Integrates vendor warranty verification (HP, Dell, Lenovo) into lifecycle decisions",
    "Supports 4 enterprise workflows: Handover, Return, Self-Declaration, Ad-hoc",
    "Full audit trail with DynamoDB history and CloudWatch logging",
    "Role-based access with Cognito authentication (Employee / Admin)",
], top=3.0)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: Problem Statement
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("Problem")
add_title_box(slide, "The Problem", top=0.3, size=28)
add_bullet_box(slide, [
    "Manual asset inspections are slow (30-45 min per device), subjective, and inconsistent",
    "No standardized damage scoring — depends on who inspects",
    "Fraud is undetectable: employees submit old photos, stock images, or photos of different devices",
    "Warranty status is checked manually on vendor websites — time-consuming and often skipped",
    "Device type mismatch goes unnoticed (claiming monitor inspection but submitting laptop photos)",
    "No single source of truth for asset condition history",
    "Lifecycle decisions (repair/replace/continue) made without data-driven justification",
    "Audit compliance gaps — no tamper-proof evidence trail",
], top=1.2, size=15)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: Solution Overview
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("Solution")
add_title_box(slide, "Solution: Multi-Agent AI Pipeline", top=0.3, size=28)
add_text_box(slide, "Six specialized AI agents execute sequentially with fail-fast logic:", top=1.0, size=15, color=WHITE)
add_bullet_box(slide, [
    "1. Identity Verification — OCR + CMDB cross-reference (auto-registers unknown assets)",
    "2. Vision Quality Gate — Ensures photos are clear, well-lit, and show the full device",
    "3. Fraud Detection — EXIF analysis, stock photo detection, liveness check, duplicate detection",
    "4. Damage Assessment — AI vision scores physical damage (0-100) with device type auto-detection",
    "5. Lifecycle Decision — Rule-based repair/replace/continue logic using damage + age + warranty",
    "6. Warranty Verification — Real-time check against HP/Dell/Lenovo vendor APIs",
    "",
    "Pipeline halts immediately on Critical fraud (stock photos) or identity failure.",
    "Results include confidence scores, evidence, and actionable recommendations.",
], top=1.8, size=14)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: Architecture Diagram
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("Architecture")
add_title_box(slide, "Technical Architecture", top=0.3, size=28)
add_text_box(slide, """
┌─────────────────────────────────────────────────────────────────────────────────┐
│  FRONTEND (CloudFront + S3)          │  BACKEND (Lambda + API Gateway)           │
│                                      │                                           │
│  Single-page portal (index.html)     │  assetguardian-harness-invoke Lambda      │
│  Cognito Auth (SigV4 signing)        │  ├── handler.py (routing + auth)          │
│  S3 direct upload (presigned)        │  ├── orchestrator.py (pipeline)           │
│  Dark theme, responsive              │  └── agents/                              │
│                                      │      ├── identity_verification.py         │
│  CloudFront Distribution             │      ├── vision_quality_gate.py           │
│  dn40ox0gzb4b1.cloudfront.net        │      ├── fraud_detection.py               │
│                                      │      ├── damage_detection.py              │
│  API CloudFront                      │      ├── lifecycle_decision.py            │
│  d2xsszrekq4050.cloudfront.net       │      ├── warranty_verification.py         │
│                                      │      ├── display_health.py                │
│                                      │      └── evaluators.py                    │
└─────────────────────────────────────────────────────────────────────────────────┘

┌─────────────────────────────────────────────────────────────────────────────────┐
│  DATA LAYER                          │  AI SERVICES                              │
│                                      │                                           │
│  DynamoDB: assetguardian-cmdb        │  Amazon Bedrock (Claude Sonnet 4)         │
│  DynamoDB: inspection-history        │  Amazon Rekognition (OCR / DetectText)    │
│  DynamoDB: warranty-cache            │  Vendor APIs (HP/Dell/Lenovo)             │
│  S3: asset-photos bucket             │                                           │
│  Cognito User Pool                   │  WAFv2 Web ACL (rate limiting)            │
│  Secrets Manager (origin verify)     │  X-Ray Tracing (performance)              │
└─────────────────────────────────────────────────────────────────────────────────┘
""", top=1.0, left=0.3, width=12.5, size=10, color=ACCENT_GREEN)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: Workflows
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("Workflows")
add_title_box(slide, "Enterprise Workflows", top=0.3, size=28)
add_text_box(slide, "Four inspection workflows with SLA targets:", top=1.0, size=15, color=WHITE)
add_bullet_box(slide, [
    "Employee Handover (24h SLA)",
    "  Identity → Vision → Fraud → Damage → Lifecycle → Warranty",
    "  Used when assigning a new device to an employee",
    "",
    "Asset Return (48h SLA)",
    "  Identity → Vision → Fraud → Damage → Display Health → Lifecycle → Warranty",
    "  Most comprehensive — includes display pixel/backlight checks",
    "",
    "Annual Self-Declaration (72h SLA)",
    "  Identity → Attestation → Vision → Fraud → Damage → Lifecycle → Warranty",
    "  Yearly condition check with employee attestation",
    "",
    "Ad-hoc Inspection (4h critical / 24h standard SLA)",
    "  Identity → Vision → Fraud → Damage (conditional) → Display (conditional) → Lifecycle → Warranty",
    "  Flexible — triggered by IT for specific concerns",
], top=1.6, size=13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: AI Agents Detail
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("AI Agents")
add_title_box(slide, "AI Agent Details", top=0.3, size=28)
add_bullet_box(slide, [
    "IDENTITY VERIFICATION",
    "  • Rekognition DetectText OCR reads asset tag serial numbers",
    "  • Cross-references DynamoDB CMDB (asset registry)",
    "  • Auto-registers unknown assets (creates CMDB entry on first inspection)",
    "  • Matches employee ID / email against assigned user",
    "",
    "FRAUD DETECTION (5-check framework)",
    "  • EXIF metadata presence and validity",
    "  • Image freshness (captured within 48 hours)",
    "  • Liveness check (optional OTP code visible in photo)",
    "  • Duplicate detection (SHA-256 hash comparison)",
    "  • AI stock photo detection (Claude Vision — rejects studio/internet images)",
    "  • Critical fraud halts the entire pipeline immediately",
    "",
    "DAMAGE DETECTION",
    "  • Auto-detects device type (laptop/monitor/phone/tablet) from photos",
    "  • Flags device type mismatch (declared vs actual)",
    "  • Scores damage 0-100 with severity category (None/Minor/Moderate/Severe/Critical)",
    "  • Identifies specific damage types: cracks, scratches, dents, missing keys, hinge damage",
], top=1.0, size=12)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: Warranty Integration
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("Warranty")
add_title_box(slide, "Vendor Warranty Verification", top=0.3, size=28)
add_text_box(slide, "Real-time warranty status from principal vendors, integrated into lifecycle decisions:", top=1.0, size=15, color=WHITE)
add_bullet_box(slide, [
    "Supported Vendors:",
    "  • HP — Enterprise Warranty API (OAuth2, API key/secret)",
    "  • Dell — TechDirect API (OAuth2 client credentials)",
    "  • Lenovo — Public warranty lookup (no credentials needed)",
    "",
    "How it works:",
    "  • Device vendor auto-detected from CMDB record (device model / manufacturer field)",
    "  • Serial number queried against vendor API with 5-second timeout",
    "  • Results cached in DynamoDB (24h TTL) to avoid repeated lookups",
    "  • Warranty status (ACTIVE/EXPIRED) feeds into lifecycle decision",
    "",
    "Data returned:",
    "  • Warranty start/end dates, days remaining",
    "  • Service level (e.g., Next Business Day Onsite)",
    "  • Coverage details (multiple entitlements)",
    "  • Product name and model confirmation from vendor records",
    "",
    "Runs in PARALLEL with damage detection — adds zero latency to the pipeline",
], top=1.6, size=13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: Security
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("Security")
add_title_box(slide, "Security Architecture", top=0.3, size=28)
add_bullet_box(slide, [
    "Authentication & Authorization:",
    "  • Amazon Cognito User Pool (email/password + verification)",
    "  • SigV4 request signing for all API calls",
    "  • AWS_IAM authorization on API Gateway routes",
    "  • Role-based access: Employee vs Admin",
    "",
    "Origin Protection:",
    "  • CloudFront origin-verify header (shared secret via Secrets Manager)",
    "  • WAFv2 Web ACL attached to API CloudFront distribution",
    "  • Direct API Gateway access blocked — must go through CloudFront",
    "",
    "Data Protection:",
    "  • Photos uploaded directly to S3 with presigned URLs (never touch server)",
    "  • S3 bucket policy restricts access to authenticated users",
    "  • DynamoDB encryption at rest (AWS managed keys)",
    "  • Inspection history has TTL-based expiration",
    "",
    "Input Sanitization:",
    "  • Prompt injection defense (sanitize.py strips all special characters)",
    "  • Asset ID, serial number, employee ID validated with strict regex",
    "  • S3 key path validation prevents traversal attacks",
], top=1.0, size=12)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: Infrastructure
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("Infrastructure")
add_title_box(slide, "AWS Infrastructure (CDK)", top=0.3, size=28)
add_bullet_box(slide, [
    "Compute:",
    "  • Lambda (Python 3.12) — 512MB, 300s timeout, x86_64",
    "  • Single function handles all routes (orchestrator pattern)",
    "",
    "Networking:",
    "  • CloudFront (portal) — dn40ox0gzb4b1.cloudfront.net",
    "  • CloudFront (API) — d2xsszrekq4050.cloudfront.net (60s origin timeout)",
    "  • API Gateway HTTP API (v2) — rwgheig8j7.execute-api",
    "",
    "Storage:",
    "  • S3: assetguardian-portal (frontend hosting)",
    "  • S3: assetguardian-asset-photos (inspection evidence)",
    "  • DynamoDB: assetguardian-cmdb (asset registry, 295 assets)",
    "  • DynamoDB: assetguardian-inspection-history (results + audit)",
    "  • DynamoDB: assetguardian-warranty-cache (vendor API cache, 24h TTL)",
    "",
    "AI Services:",
    "  • Amazon Bedrock — Claude Haiku 4.5 (vision + reasoning)",
    "  • Amazon Rekognition — DetectText (OCR for asset tags)",
    "",
    "Deployment: AWS CDK (Python) — single stack, single region (ap-southeast-1)",
], top=1.0, size=12)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11: Results & Reporting
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("Results")
add_title_box(slide, "Results & Reporting", top=0.3, size=28)
add_bullet_box(slide, [
    "Inspection Results (displayed on-screen immediately after pipeline completes):",
    "  • Identity Verification — verified/not, CMDB status, serial match, ownership",
    "  • Vision Quality — overall score, sharpness, lighting",
    "  • Fraud Detection — risk level, legitimacy confidence, stock photo flag",
    "  • Damage Assessment — severity score, condition grade, detected device type/brand",
    "  • Lifecycle Decision — continue/repair/replace, cost estimate, rule applied",
    "  • Warranty Verification — status, end date, days remaining, service level",
    "",
    "Actions available after inspection:",
    "  • 📄 Save as PDF — opens print-friendly report in new tab",
    "  • 📧 Email to Employee — pre-filled email with summary",
    "  • 🔄 New Inspection — clears form and starts fresh",
    "",
    "Audit & Compliance:",
    "  • Full history stored in DynamoDB with TTL",
    "  • Event Logs tab (admin only) shows CloudWatch logs",
    "  • Each inspection tagged with reference ID for traceability",
    "  • Evidence photos retained in S3 with audit prefix",
], top=1.0, size=13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12: Demo Scenario
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("Demo")
add_title_box(slide, "Demo Scenario", top=0.3, size=28)
add_text_box(slide, "Scenario: Healthy HP EliteBook — Employee Handover", top=1.0, size=16, color=ACCENT_GREEN)
add_bullet_box(slide, [
    "Portal URL: https://dn40ox0gzb4b1.cloudfront.net",
    "",
    "Steps:",
    "  1. Sign in with Cognito credentials",
    "  2. Select workflow: Employee Handover",
    "  3. Enter Asset ID: ASSET-0001, Serial: 5CG01523C7, Employee: EMP1001",
    "  4. Upload device photos (front, rear, side) — real workplace photos",
    "  5. Upload asset tag photo",
    "  6. Click 'Run Inspection'",
    "",
    "Expected Result:",
    "  • Identity: Verified (matched in CMDB)",
    "  • Vision: Passed",
    "  • Fraud: Low risk (genuine workplace photos)",
    "  • Damage: Minor (score < 15), condition Good",
    "  • Lifecycle: Continue Use ($0 cost)",
    "  • Warranty: Active (HP, Next Business Day)",
    "",
    "Negative Test: Upload stock photo → Pipeline halts at Fraud Detection (Critical)",
], top=1.6, size=13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13: Roadmap
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("Roadmap")
add_title_box(slide, "Future Roadmap", top=0.3, size=28)
add_bullet_box(slide, [
    "Phase 2 — Planned Enhancements:",
    "  • Mobile-native app (React Native) with offline inspection capability",
    "  • Barcode/QR code scanning for instant asset lookup",
    "  • Automated email notifications on inspection completion (SES integration)",
    "  • Bulk inspection via CSV upload",
    "  • Manager approval workflow for high-value lifecycle decisions",
    "  • Integration with ServiceNow / Jira for repair ticket creation",
    "",
    "Phase 3 — Advanced AI:",
    "  • Predictive failure analysis based on inspection history trends",
    "  • Automated procurement recommendations when fleet warranty expires",
    "  • Multi-language support for global deployments",
    "  • Video-based inspection (real-time damage detection from walk-around clip)",
    "  • Bedrock AgentCore Gateway integration for conversational inspection",
    "",
    "Phase 4 — Enterprise Scale:",
    "  • Multi-tenant SaaS deployment",
    "  • Custom damage scoring models per organization",
    "  • API marketplace for third-party integrations",
], top=1.0, size=13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14: Thank You
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("End")
add_title_box(slide, "Thank You", top=2.5, left=1, size=40, color=ACCENT_GREEN)
add_text_box(slide, "AssetGuardian AI — Making Enterprise Asset Management Intelligent", top=3.5, left=1, size=18, color=MUTED)
add_text_box(slide, "Portal: https://dn40ox0gzb4b1.cloudfront.net\nGitHub: https://github.com/harishov/assetguardian", top=4.5, left=1, size=14, color=ACCENT_BLUE)
add_text_box(slide, "NCS Technology Services  |  Built with Amazon Bedrock, Lambda, CloudFront, DynamoDB", top=6.0, left=1, size=12, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "docs", "AssetGuardian_Solution_Document.pptx")
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f"Created: {output_path}")
